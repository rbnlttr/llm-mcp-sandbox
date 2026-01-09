from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Optional, Dict
import os
import anthropic
from pathlib import Path
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import io
import httpx
import json
from datetime import datetime

app = FastAPI(title="LLM MCP Sandbox API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = Path(os.getenv("UPLOAD_DIR", "./uploads"))
PROJECT_DIR = Path(os.getenv("PROJECT_DIR", "./project"))
REFERENCE_DIR = Path(os.getenv("REFERENCE_DIR", "./reference"))

UPLOAD_DIR.mkdir(exist_ok=True)
PROJECT_DIR.mkdir(exist_ok=True)
REFERENCE_DIR.mkdir(exist_ok=True)

# Configuration
USE_LOCAL_LLM = os.getenv("USE_LOCAL_LLM", "true").lower() == "true"
OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://ollama:11434")
LOCAL_MODEL = os.getenv("LOCAL_MODEL", "llama3.2")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
ANTHROPIC_AVAILABLE = ANTHROPIC_API_KEY and ANTHROPIC_API_KEY.strip()

client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY) if (ANTHROPIC_AVAILABLE) else None
# Cache for directory contents
_directory_cache = {}


class ChatRequest(BaseModel):
    message: str
    documents: List[dict]
    use_local: Optional[bool] = None
    include_project: Optional[bool] = True
    include_reference: Optional[bool] = True


class ModelInfo(BaseModel):
    name: str
    available: bool
    type: str


class DirectoryInfo(BaseModel):
    path: str
    files: List[dict]
    total_size: int
    file_count: int


def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from PDF"""
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF extraction failed: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    """Extract text from DOCX"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"DOCX extraction failed: {str(e)}")


def extract_text_from_xlsx(file_content: bytes) -> str:
    """Extract text from XLSX"""
    try:
        xlsx_file = io.BytesIO(file_content)
        workbook = openpyxl.load_workbook(xlsx_file)
        text = ""
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"\n=== Sheet: {sheet_name} ===\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"XLSX extraction failed: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from PPTX"""
    try:
        pptx_file = io.BytesIO(file_content)
        presentation = Presentation(pptx_file)
        text = ""
        for i, slide in enumerate(presentation.slides, 1):
            text += f"\n=== Slide {i} ===\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PPTX extraction failed: {str(e)}")


def extract_text_from_file(file_path: Path) -> str:
    """Extract text from any supported file"""
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
        
        filename = file_path.name.lower()
        
        if filename.endswith('.pdf'):
            return extract_text_from_pdf(content)
        elif filename.endswith('.docx'):
            return extract_text_from_docx(content)
        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            return extract_text_from_xlsx(content)
        elif filename.endswith('.pptx'):
            return extract_text_from_pptx(content)
        elif filename.endswith('.txt') or filename.endswith('.md'):
            return content.decode('utf-8')
        elif filename.endswith(('.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.go', '.rs')):
            return content.decode('utf-8')
        else:
            return f"[Unsupported file type: {filename}]"
    except Exception as e:
        return f"[Error reading {file_path.name}: {str(e)}]"


def scan_directory(directory: Path, max_files: int = 100) -> Dict:
    """Scan directory and extract content from all supported files"""
    if not directory.exists():
        return {"files": [], "total_size": 0, "file_count": 0}
    
    files = []
    total_size = 0
    file_count = 0
    
    supported_extensions = {
        '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx',
        '.txt', '.md', '.py', '.js', '.java', '.cpp', '.c', 
        '.h', '.cs', '.go', '.rs', '.json', '.yaml', '.yml'
    }
    
    for file_path in directory.rglob('*'):
        if file_path.is_file() and file_count < max_files:
            if file_path.suffix.lower() in supported_extensions:
                try:
                    size = file_path.stat().st_size
                    content = extract_text_from_file(file_path)
                    
                    relative_path = file_path.relative_to(directory)
                    
                    files.append({
                        "name": file_path.name,
                        "path": str(relative_path),
                        "size": size,
                        "content": content,
                        "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
                    })
                    
                    total_size += size
                    file_count += 1
                except Exception as e:
                    print(f"Error processing {file_path}: {e}")
    
    return {
        "files": files,
        "total_size": total_size,
        "file_count": file_count
    }


async def chat_with_ollama(prompt: str, system_prompt: str) -> dict:
    """Chat with local Ollama LLM"""
    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            response = await client.post(
                f"{OLLAMA_HOST}/api/generate",
                json={
                    "model": LOCAL_MODEL,
                    "prompt": f"System: {system_prompt}\n\nUser: {prompt}",
                    "stream": False,
                    "options": {
                        "temperature": 0.7,
                        "top_p": 0.9,
                    }
                }
            )
            
            if response.status_code != 200:
                raise HTTPException(status_code=500, detail=f"Ollama error: {response.text}")
            
            result = response.json()
            return {
                "response": result.get("response", ""),
                "model": LOCAL_MODEL,
                "usage": {
                    "prompt_tokens": result.get("prompt_eval_count", 0),
                    "completion_tokens": result.get("eval_count", 0),
                    "total_tokens": result.get("prompt_eval_count", 0) + result.get("eval_count", 0)
                }
            }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Ollama chat failed: {str(e)}")


def chat_with_claude(prompt: str, system_prompt: str) -> dict:
    """Chat with Claude API"""
    if not client:
        raise HTTPException(status_code=500, detail="Anthropic API key not configured")
    
    try:
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4000,
            system=system_prompt,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        
        response_text = ""
        for block in message.content:
            if block.type == "text":
                response_text += block.text
        
        return {
            "response": response_text,
            "model": message.model,
            "usage": {
                "input_tokens": message.usage.input_tokens,
                "output_tokens": message.usage.output_tokens,
                "total_tokens": message.usage.input_tokens + message.usage.output_tokens
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Claude chat failed: {str(e)}")


@app.get("/")
async def root():
    return {
        "status": "ok",
        "message": "LLM MCP Sandbox API",
        "local_llm_enabled": USE_LOCAL_LLM,
        "local_model": LOCAL_MODEL if USE_LOCAL_LLM else None,
        "directories": {
            "project": str(PROJECT_DIR),
            "reference": str(REFERENCE_DIR),
            "uploads": str(UPLOAD_DIR)
        }
    }


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload and process a document"""
    try:
        content = await file.read()
        filename = file.filename.lower()
        
        # Extract text based on file type
        if filename.endswith('.pdf'):
            text = extract_text_from_pdf(content)
        elif filename.endswith('.docx'):
            text = extract_text_from_docx(content)
        elif filename.endswith('.doc'):
            raise HTTPException(status_code=400, detail="DOC format not supported, use DOCX")
        elif filename.endswith('.xlsx') or filename.endswith('.xls'):
            text = extract_text_from_xlsx(content)
        elif filename.endswith('.pptx'):
            text = extract_text_from_pptx(content)
        elif filename.endswith('.txt'):
            text = content.decode('utf-8')
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format")
        
        return {
            "filename": file.filename,
            "content": text,
            "size": len(content),
            "status": "processed"
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/directories/project")
async def get_project_directory():
    """Get project directory contents"""
    result = scan_directory(PROJECT_DIR)
    return {
        "path": str(PROJECT_DIR),
        "files": result["files"],
        "total_size": result["total_size"],
        "file_count": result["file_count"]
    }


@app.get("/directories/reference")
async def get_reference_directory():
    """Get reference directory contents"""
    result = scan_directory(REFERENCE_DIR)
    return {
        "path": str(REFERENCE_DIR),
        "files": result["files"],
        "total_size": result["total_size"],
        "file_count": result["file_count"]
    }


@app.post("/directories/refresh")
async def refresh_directories():
    """Manually refresh directory cache"""
    project = scan_directory(PROJECT_DIR)
    reference = scan_directory(REFERENCE_DIR)
    
    return {
        "project": {
            "file_count": project["file_count"],
            "total_size": project["total_size"]
        },
        "reference": {
            "file_count": reference["file_count"],
            "total_size": reference["total_size"]
        }
    }


@app.post("/chat")
async def chat(request: ChatRequest):
    """Send message with document context to LLM"""
    use_local = request.use_local if request.use_local is not None else USE_LOCAL_LLM
    
    try:
        # Build context from uploaded documents
        context = ""
        
        if request.documents:
            context += "=== HOCHGELADENE DOKUMENTE ===\n"
            for i, doc in enumerate(request.documents, 1):
                context += f"\n{i}. {doc['name']}\n"
            
            for doc in request.documents:
                context += f"\n--- {doc['name']} ---\n{doc['content']}\n"
        
        # Add project directory context
        if request.include_project:
            project_data = scan_directory(PROJECT_DIR)
            if project_data["files"]:
                context += "\n\n=== PROJEKT-DATEIEN ===\n"
                for file in project_data["files"]:
                    context += f"\n--- {file['path']} ---\n{file['content']}\n"
        
        # Add reference directory context
        if request.include_reference:
            reference_data = scan_directory(REFERENCE_DIR)
            if reference_data["files"]:
                context += "\n\n=== REFERENZDOKUMENTE (Normen/Richtlinien) ===\n"
                for file in reference_data["files"]:
                    context += f"\n--- {file['path']} ---\n{file['content']}\n"
        
        system_prompt = f"""Du bist ein AI-Assistent in einer Sandbox-Umgebung mit MCP (Model Context Protocol) für die Prüfung und Erstellung von Sicherheitsnachweisdokumenttion nach EN 50129.

{context}

Beantworte Fragen basierend auf den verfügbaren Dokumenten im Kontext der Sicherheitsnachweisführung. Sei präzise und zitiere die Quelle (inkl. Dateipfad). 
Wenn du auf Normen oder Richtlinien Bezug nimmst, nenne sie explizit."""
        
        # Choose LLM
        if use_local:
            result = await chat_with_ollama(request.message, system_prompt)
            result["llm_type"] = "local"
        else:
            result = chat_with_claude(request.message, system_prompt)
            result["llm_type"] = "claude"
        
        return result
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/models")
async def get_models():
    """Get available models"""
    models = []
    
    # Check local Ollama
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            response = await client.get(f"{OLLAMA_HOST}/api/tags")
            if response.status_code == 200:
                ollama_models = response.json().get("models", [])
                for model in ollama_models:
                    models.append({
                        "name": model["name"],
                        "available": True,
                        "type": "local",
                        "size": model.get("size", 0)
                    })
    except:
        models.append({
            "name": LOCAL_MODEL,
            "available": False,
            "type": "local",
            "error": "Ollama not available"
        })
    
    # Check Claude
    models.append({
        "name": "claude-sonnet-4",
        "available": ANTHROPIC_AVAILABLE,
        "type": "cloud"
    })
    
    return {"models": models}


@app.get("/health")
async def health():
    """Health check"""
    ollama_available = False
    try:
        async with httpx.AsyncClient(timeout=5.0) as client:
            response = await client.get(f"{OLLAMA_HOST}/api/tags")
            ollama_available = response.status_code == 200
    except:
        pass
    
    # Check directories
    project_exists = PROJECT_DIR.exists()
    reference_exists = REFERENCE_DIR.exists()
    
    project_files = len(list(PROJECT_DIR.rglob('*'))) if project_exists else 0
    reference_files = len(list(REFERENCE_DIR.rglob('*'))) if reference_exists else 0
    
    return {
        "status": "healthy",
        "ollama_available": ollama_available,
        "ollama_host": OLLAMA_HOST,
        "claude_available": ANTHROPIC_AVAILABLE,
        "default_llm": "local" if USE_LOCAL_LLM else "cloud",
        "directories": {
            "project": {
                "exists": project_exists,
                "files": project_files,
                "path": str(PROJECT_DIR)
            },
            "reference": {
                "exists": reference_exists,
                "files": reference_files,
                "path": str(REFERENCE_DIR)
            }
        }
    }


@app.post("/ollama/pull")
async def pull_model(model_name: str = LOCAL_MODEL):
    """Pull a model from Ollama registry"""
    try:
        async with httpx.AsyncClient(timeout=600.0) as client:
            response = await client.post(
                f"{OLLAMA_HOST}/api/pull",
                json={"name": model_name}
            )
            
            if response.status_code == 200:
                return {"status": "success", "message": f"Model {model_name} pulled successfully"}
            else:
                raise HTTPException(status_code=500, detail="Failed to pull model")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
