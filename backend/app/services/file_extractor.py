from pathlib import Path
import io, PyPDF2, docx, openpyxl
from pptx import Presentation

class FileExtractor:

    def extract(self, path: Path) -> str:
        content = path.read_bytes()
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            return self._pdf(content)
        if suffix == ".docx":
            return self._docx(content)
        if suffix in (".xlsx", ".xls"):
            return self._xlsx(content)
        if suffix == ".pptx":
            return self._pptx(content)
        if suffix in (".txt", ".md", ".py", ".json"):
            return content.decode("utf-8")

        return f"[Unsupported file type: {path.name}]"

    def _pdf(self, content: bytes) -> str:
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        return "\n".join(p.extract_text() or "" for p in reader.pages)

    def _docx(self, content: bytes) -> str:
        doc = docx.Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)

    def _xlsx(self, content: bytes) -> str:
        wb = openpyxl.load_workbook(io.BytesIO(content))
        lines = []
        for sheet in wb:
            lines.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                lines.append("\t".join("" if c is None else str(c) for c in row))
        return "\n".join(lines)

    def _pptx(self, content: bytes) -> str:
        pres = Presentation(io.BytesIO(content))
        return "\n".join(
            shape.text
            for slide in pres.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )
