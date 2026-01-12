import io
import PyPDF2
import docx
import openpyxl
from pptx import Presentation


def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from PDF content."""
    reader = PyPDF2.PdfReader(io.BytesIO(content))
    return "\n".join(p.extract_text() or "" for p in reader.pages)


def extract_text_from_docx(content: bytes) -> str:
    """Extract text from DOCX content."""
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_text_from_xlsx(content: bytes) -> str:
    """Extract text from XLSX content."""
    wb = openpyxl.load_workbook(io.BytesIO(content))
    lines = []
    for sheet in wb:
        lines.append(f"=== Sheet: {sheet.title} ===")
        for row in sheet.iter_rows(values_only=True):
            lines.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(lines)


def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from PPTX content."""
    pres = Presentation(io.BytesIO(content))
    return "\n".join(
        shape.text
        for slide in pres.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )